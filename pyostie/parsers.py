import docx2txt
import xlrd
import csv
import cv2
import pytesseract
from PIL import Image
import PyPDF2
import pdfplumber
from pptx import Presentation
from pdf2image import convert_from_path
from pyostie.convert import *
from pyostie.insights_ext import *

a = pd.DataFrame()


class DOCXParser:

    def __init__(self, filename):
        """

        :param filename:
        """
        self.file = filename

    def extract_docx(self):
        """

        :return:
        """
        output = docx2txt.process(self.file)
        return output


class XLSXParser:

    def __init__(self, filename):
        """

        :param filename:
        """
        self.file = filename

    def extract_xlsx(self):
        """

        :return:
        """
        out_list = []
        book = xlrd.open_workbook(self.file)
        for val in range(len(book.sheet_names())):
            sheet = book.sheet_by_index(val)
            for res in range(sheet.nrows):
                output = " " + " ".join(str(val_) for val_ in (sheet.row_values(res)))
                out_list.append(output)
        return out_list


class CSVParser:

    def __init__(self, filename, delimiter=','):
        """

        :param filename:
        """
        self.file = filename
        self.delimiter = delimiter

    def extract_csv(self):
        """

        :return:
        """
        with open(self.file) as file:
            output = csv.reader(file, delimiter=self.delimiter)
            return ' '.join([' '.join(row) for row in output])


class ImageParser:

    def __init__(self, filename, tess_path=None):
        """

        :param filename:
        :param tess_path
        """
        self.file = filename
        self.path = tess_path

    def extract_image(self):
        """

        :return:
        """
        out_list = []
        if self.path is not None:
            pytesseract.pytesseract.tesseract_cmd = self.path
            img = Image.open(self.file)
            text = pytesseract.image_to_string(img)
            out_list.append(text)
        else:
            img = Image.open(self.file)
            text = pytesseract.image_to_string(img)
            out_list.append(text)
        return out_list


class PDFParser:
    """

    """

    def __init__(self, filename, insights=False):
        """

        :param filename:
        :param insights:
        """
        self.file = filename
        self.insights = insights

    def extract_pypdf2(self):
        """

        :param kwargs:
        :return:
        """

        global __insights
        contents = []
        text = ' '
        pdfFileObj = open(self.file, 'rb')
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        pdfPages = pdfReader.getNumPages()
        if pdfPages == 1:
            for val in range(pdfReader.numPages):
                pageObject = pdfReader.getPage(val)
                text = text + pageObject.extractText()
            contents.append(text)
            if self.insights:
                conv = conversion(self.file)
                __conv = conv.convert()
                insights = generate_insights(__conv, df)
                __insights = insights.generate_df()
                remove_files(__conv)
                return __insights, contents
            else:
                return contents

        if pdfPages >= 2:
            print("We are currently not generating insights for multi page pdf's, so returning only text."
                  " Please watch this space for updates.")
            contents = []
            text = ' '
            pdfFileObj = open(self.file, 'rb')
            pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
            pdfPages = pdfReader.getNumPages()
            for val in range(pdfPages):
                pageObject = pdfReader.getPage(val)
                text = text + pageObject.extractText()
            contents.append(text)
            emp_df = pd.DataFrame()
            return emp_df, contents

    def extract_pdfplumber(self):
        """

        :return:
        """
        out_list = []
        with pdfplumber.open(self.file) as pdf:
            for val in range(len(pdf.pages)):
                page = pdf.pages[val]
                output = page.extract_text()
                out_list.append(output)
        return out_list


class TXTParser:

    def __init__(self, filename):
        """

        :param filename:
        """
        self.file = filename

    def extract_txt(self):
        """

        :return:
        """
        with open(self.file) as file:
            return file.read()


class PPTXParser:

    def __init__(self, filename):
        """

        :param filename:
        """
        self.file = filename

    def extract_pptx(self):
        """
        :return:
        """
        text = []
        paper = Presentation(self.file)
        for slide in paper.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    stripped = paragraph.text.strip()
                    if stripped:
                        text.append(paragraph.text)
        return text
